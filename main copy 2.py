from fastapi import FastAPI, UploadFile, File, HTTPException, Query
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional, Dict, Any, List
from io import BytesIO
import os
import logging
import csv
import json
import re
import unicodedata
from datetime import datetime

# Core libraries
from PyPDF2 import PdfReader
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
from bs4 import BeautifulSoup
import markdown
import openpyxl
import xlrd
from ebooklib import epub
from striprtf.striprtf import rtf_to_text
from tempfile import NamedTemporaryFile

# Data processing libraries
import pandas as pd
import numpy as np
from textblob import TextBlob
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize, sent_tokenize
from nltk.stem import PorterStemmer, WordNetLemmatizer

app = FastAPI(
    title="Advanced File Parser & Data Processor",
    description="Extract, clean, and preprocess text from various file formats",
    version="2.0.0"
)

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('parser.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("parser")

# Download required NLTK data
try:
    nltk.data.find('tokenizers/punkt')
    nltk.data.find('corpora/stopwords')
    nltk.data.find('corpora/wordnet')
except LookupError:
    nltk.download('punkt')
    nltk.download('stopwords')
    nltk.download('wordnet')
    nltk.download('punkt_tab')

# Pydantic models
class ProcessingOptions(BaseModel):
    remove_extra_spaces: bool = True
    remove_special_chars: bool = False
    normalize_unicode: bool = True
    remove_stopwords: bool = False
    lowercase: bool = False
    remove_numbers: bool = False
    remove_punctuation: bool = False
    min_word_length: int = 1
    max_text_length: int = 50000
    language: str = "english"
    stemming: bool = False
    lemmatization: bool = False
    remove_urls: bool = True
    remove_emails: bool = True
    remove_extra_newlines: bool = True
    preserve_structure: bool = True

class ProcessingResult(BaseModel):
    text: str
    metadata: Dict[str, Any]
    statistics: Dict[str, Any]
    processing_time: float
    file_info: Dict[str, Any]

# --- Data Cleaning and Preprocessing ---

class TextCleaner:
    def __init__(self, options: ProcessingOptions):
        self.options = options
        self.stemmer = PorterStemmer()
        self.lemmatizer = WordNetLemmatizer()
        self.stop_words = set(stopwords.words(options.language))
    
    def clean_text(self, text: str) -> str:
        """Comprehensive text cleaning pipeline"""
        if not text:
            return ""
        
        original_length = len(text)
        
        # Unicode normalization
        if self.options.normalize_unicode:
            text = unicodedata.normalize('NFKD', text)
        
        # Remove URLs
        if self.options.remove_urls:
            text = re.sub(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+', '', text)
        
        # Remove email addresses
        if self.options.remove_emails:
            text = re.sub(r'\S+@\S+', '', text)
        
        # Remove extra newlines and spaces
        if self.options.remove_extra_newlines:
            text = re.sub(r'\n\s*\n', '\n\n', text)  # Multiple newlines to double newline
        
        if self.options.remove_extra_spaces:
            text = re.sub(r'\s+', ' ', text)  # Multiple spaces to single space
            text = text.strip()
        
        # Remove special characters (keep basic punctuation if preserve_structure is True)
        if self.options.remove_special_chars:
            if self.options.preserve_structure:
                text = re.sub(r'[^\w\s.,!?;:()\-"\']', '', text)
            else:
                text = re.sub(r'[^\w\s]', '', text)
        
        # Convert to lowercase
        if self.options.lowercase:
            text = text.lower()
        
        # Remove numbers
        if self.options.remove_numbers:
            text = re.sub(r'\d+', '', text)
        
        # Remove punctuation
        if self.options.remove_punctuation:
            text = re.sub(r'[^\w\s]', '', text)
        
        # Advanced NLP processing
        if any([self.options.remove_stopwords, self.options.stemming, self.options.lemmatization]):
            text = self._advanced_processing(text)
        
        # Filter by minimum word length
        if self.options.min_word_length > 1:
            words = text.split()
            words = [word for word in words if len(word) >= self.options.min_word_length]
            text = ' '.join(words)
        
        # Truncate if necessary
        if len(text) > self.options.max_text_length:
            text = text[:self.options.max_text_length] + "..."
        
        return text.strip()
    
    def _advanced_processing(self, text: str) -> str:
        """Advanced NLP processing with tokenization"""
        try:
            sentences = sent_tokenize(text)
            processed_sentences = []
            
            for sentence in sentences:
                tokens = word_tokenize(sentence)
                
                # Remove stopwords
                if self.options.remove_stopwords:
                    tokens = [token for token in tokens if token.lower() not in self.stop_words]
                
                # Stemming
                if self.options.stemming:
                    tokens = [self.stemmer.stem(token) for token in tokens]
                
                # Lemmatization
                if self.options.lemmatization:
                    tokens = [self.lemmatizer.lemmatize(token) for token in tokens]
                
                processed_sentences.append(' '.join(tokens))
            
            return ' '.join(processed_sentences)
        except Exception as e:
            logger.warning(f"Advanced processing failed: {e}, returning original text")
            return text
    
    def get_statistics(self, original_text: str, cleaned_text: str) -> Dict[str, Any]:
        """Generate text statistics"""
        try:
            blob_original = TextBlob(original_text)
            blob_cleaned = TextBlob(cleaned_text)
            
            return {
                "original_length": len(original_text),
                "cleaned_length": len(cleaned_text),
                "compression_ratio": round(len(cleaned_text) / len(original_text), 3) if original_text else 0,
                "original_words": len(blob_original.words),
                "cleaned_words": len(blob_cleaned.words),
                "original_sentences": len(blob_original.sentences),
                "cleaned_sentences": len(blob_cleaned.sentences),
                "detected_language": str(blob_cleaned.detect_language()) if cleaned_text else "unknown"
            }
        except Exception as e:
            logger.warning(f"Statistics generation failed: {e}")
            return {
                "original_length": len(original_text),
                "cleaned_length": len(cleaned_text),
                "compression_ratio": round(len(cleaned_text) / len(original_text), 3) if original_text else 0,
                "error": str(e)
            }

# --- Enhanced Text Extractors ---

def extract_pdf_text(file_bytes: bytes) -> tuple[str, dict]:
    """Extract text from PDF with metadata"""
    metadata = {"pages": 0, "extraction_method": ""}
    
    try:
        # Try PyMuPDF first (better performance)
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        text = "\n".join(page.get_text() for page in doc)
        metadata["pages"] = doc.page_count
        metadata["extraction_method"] = "PyMuPDF"
        doc.close()
    except Exception as e:
        logger.warning("PyMuPDF failed, fallback to PyPDF2: %s", e)
        try:
            reader = PdfReader(BytesIO(file_bytes))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            metadata["pages"] = len(reader.pages)
            metadata["extraction_method"] = "PyPDF2"
        except Exception as e2:
            logger.error("Both PDF extractors failed: %s", e2)
            raise HTTPException(status_code=500, detail=f"PDF extraction failed: {e2}")
    
    return text.strip(), metadata

def extract_docx_text(file_bytes: bytes) -> tuple[str, dict]:
    """Extract text from DOCX with metadata"""
    doc = Document(BytesIO(file_bytes))
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    
    # Extract tables
    table_text = []
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells]
            table_text.append(" | ".join(row_text))
    
    text = "\n".join(paragraphs)
    if table_text:
        text += "\n\nTables:\n" + "\n".join(table_text)
    
    metadata = {
        "paragraphs": len(paragraphs),
        "tables": len(doc.tables),
        "extraction_method": "python-docx"
    }
    
    return text.strip(), metadata

def extract_pptx_text(file_bytes: bytes) -> tuple[str, dict]:
    """Extract text from PPTX with metadata"""
    prs = Presentation(BytesIO(file_bytes))
    slides_text = []
    
    for i, slide in enumerate(prs.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())
        
        if slide_content:
            slides_text.append(f"Slide {i}:\n" + "\n".join(slide_content))
    
    metadata = {
        "slides": len(prs.slides),
        "slides_with_text": len(slides_text),
        "extraction_method": "python-pptx"
    }
    
    return "\n\n".join(slides_text), metadata

def extract_image_text(file_bytes: bytes) -> tuple[str, dict]:
    """Extract text from images using OCR"""
    image = Image.open(BytesIO(file_bytes))
    
    # Convert to RGB if necessary
    if image.mode != 'RGB':
        image = image.convert('RGB')
    
    text = pytesseract.image_to_string(image)
    
    metadata = {
        "image_size": image.size,
        "image_mode": image.mode,
        "extraction_method": "pytesseract"
    }
    
    return text.strip(), metadata

def extract_csv_text(file_bytes: bytes) -> tuple[str, dict]:
    """Extract and structure CSV data"""
    try:
        # Try pandas first for better handling
        df = pd.read_csv(BytesIO(file_bytes))
        text = df.to_string(index=False)
        
        metadata = {
            "rows": len(df),
            "columns": len(df.columns),
            "column_names": list(df.columns),
            "extraction_method": "pandas"
        }
    except Exception as e:
        logger.warning(f"Pandas CSV read failed: {e}, falling back to csv module")
        text_content = file_bytes.decode("utf-8", errors="ignore")
        reader = csv.reader(text_content.splitlines())
        rows = list(reader)
        text = "\n".join([", ".join(row) for row in rows])
        
        metadata = {
            "rows": len(rows),
            "extraction_method": "csv_module",
            "error": str(e)
        }
    
    return text, metadata

def extract_excel_text(file_bytes: bytes, ext: str) -> tuple[str, dict]:
    """Extract text from Excel files"""
    with NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        if ext == ".xlsx":
            df_dict = pd.read_excel(tmp_path, sheet_name=None)
            sheets_text = []
            
            for sheet_name, df in df_dict.items():
                sheet_text = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                sheets_text.append(sheet_text)
            
            text = "\n\n".join(sheets_text)
            metadata = {
                "sheets": len(df_dict),
                "sheet_names": list(df_dict.keys()),
                "total_rows": sum(len(df) for df in df_dict.values()),
                "extraction_method": "pandas"
            }
        else:  # .xls
            book = xlrd.open_workbook(tmp_path)
            sheets_text = []
            
            for sheet in book.sheets():
                sheet_data = []
                for row_idx in range(sheet.nrows):
                    row = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                    sheet_data.append(", ".join(row))
                sheets_text.append(f"Sheet: {sheet.name}\n" + "\n".join(sheet_data))
            
            text = "\n\n".join(sheets_text)
            metadata = {
                "sheets": len(book.sheets()),
                "extraction_method": "xlrd"
            }
            
    finally:
        os.remove(tmp_path)

    return text, metadata

# Enhanced extractors for other formats
def extract_txt_text(file_bytes: bytes) -> tuple[str, dict]:
    text = file_bytes.decode("utf-8", errors="ignore").strip()
    metadata = {"encoding": "utf-8", "extraction_method": "direct"}
    return text, metadata

def extract_html_text(file_bytes: bytes) -> tuple[str, dict]:
    soup = BeautifulSoup(file_bytes, "html.parser")
    text = soup.get_text(separator='\n').strip()
    
    metadata = {
        "title": soup.title.string if soup.title else "",
        "links": len(soup.find_all('a')),
        "images": len(soup.find_all('img')),
        "extraction_method": "BeautifulSoup"
    }
    return text, metadata

def extract_json_text(file_bytes: bytes) -> tuple[str, dict]:
    try:
        data = json.loads(file_bytes.decode("utf-8", errors="ignore"))
        text = json.dumps(data, indent=2, ensure_ascii=False)
        
        metadata = {
            "json_valid": True,
            "extraction_method": "json_module"
        }
    except json.JSONDecodeError as e:
        text = "Invalid JSON file."
        metadata = {
            "json_valid": False,
            "error": str(e),
            "extraction_method": "json_module"
        }
    
    return text, metadata

# File type mapping
EXTRACTORS = {
    ".pdf": extract_pdf_text,
    ".docx": extract_docx_text,
    ".pptx": extract_pptx_text,
    ".png": extract_image_text,
    ".jpg": extract_image_text,
    ".jpeg": extract_image_text,
    ".bmp": extract_image_text,
    ".tiff": extract_image_text,
    ".webp": extract_image_text,
    ".txt": extract_txt_text,
    ".html": extract_html_text,
    ".htm": extract_html_text,
    ".csv": extract_csv_text,
    ".json": extract_json_text,
    ".xlsx": extract_excel_text,
    ".xls": extract_excel_text,
}

# --- API Routes ---

@app.get("/")
async def root():
    return {
        "message": "Advanced File Parser & Data Processor API",
        "version": "2.0.0",
        "supported_formats": list(EXTRACTORS.keys())
    }

@app.get("/health")
async def health_check():
    return {"status": "healthy", "timestamp": datetime.now().isoformat()}

@app.post("/parse", response_model=ProcessingResult)
async def parse_file(
    file: UploadFile = File(...),
    remove_extra_spaces: bool = Query(True),
    remove_special_chars: bool = Query(False),
    normalize_unicode: bool = Query(True),
    remove_stopwords: bool = Query(False),
    lowercase: bool = Query(False),
    remove_numbers: bool = Query(False),
    remove_punctuation: bool = Query(False),
    min_word_length: int = Query(1),
    max_text_length: int = Query(50000),
    language: str = Query("english"),
    stemming: bool = Query(False),
    lemmatization: bool = Query(False),
    remove_urls: bool = Query(True),
    remove_emails: bool = Query(True)
):
    """Parse and process uploaded file with comprehensive cleaning options"""
    start_time = datetime.now()
    
    try:
        content = await file.read()
        ext = os.path.splitext(file.filename)[-1].lower()
        
        if ext not in EXTRACTORS:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type: {ext}. Supported formats: {list(EXTRACTORS.keys())}"
            )
        
        # Extract text and metadata
        if ext in [".xlsx", ".xls"]:
            raw_text, extraction_metadata = EXTRACTORS[ext](content, ext)
        else:
            raw_text, extraction_metadata = EXTRACTORS[ext](content)
        
        # Set up processing options
        options = ProcessingOptions(
            remove_extra_spaces=remove_extra_spaces,
            remove_special_chars=remove_special_chars,
            normalize_unicode=normalize_unicode,
            remove_stopwords=remove_stopwords,
            lowercase=lowercase,
            remove_numbers=remove_numbers,
            remove_punctuation=remove_punctuation,
            min_word_length=min_word_length,
            max_text_length=max_text_length,
            language=language,
            stemming=stemming,
            lemmatization=lemmatization,
            remove_urls=remove_urls,
            remove_emails=remove_emails
        )
        
        # Clean and process text
        cleaner = TextCleaner(options)
        cleaned_text = cleaner.clean_text(raw_text)
        
        # Generate statistics
        statistics = cleaner.get_statistics(raw_text, cleaned_text)
        
        processing_time = (datetime.now() - start_time).total_seconds()
        
        return ProcessingResult(
            text=cleaned_text,
            metadata=extraction_metadata,
            statistics=statistics,
            processing_time=processing_time,
            file_info={
                "filename": file.filename,
                "file_size": len(content),
                "file_type": ext,
                "content_type": file.content_type
            }
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.exception("Failed to parse file")
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")

@app.post("/batch-parse")
async def batch_parse_files(files: List[UploadFile] = File(...)):
    """Parse multiple files at once"""
    results = []
    
    for file in files:
        try:
            # Use default processing options for batch
            result = await parse_file(file)
            results.append({
                "filename": file.filename,
                "status": "success",
                "result": result
            })
        except Exception as e:
            results.append({
                "filename": file.filename,
                "status": "error",
                "error": str(e)
            })
    
    return {"results": results}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)