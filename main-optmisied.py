from fastapi import FastAPI, UploadFile, File, HTTPException, Query
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Dict, Any, List, Optional
from io import BytesIO
import os, re, unicodedata, json, csv, email, logging
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed
import asyncio
from functools import partial
from enum import Enum

# Core extraction libraries
from PyPDF2 import PdfReader
import fitz
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import cv2
from bs4 import BeautifulSoup
import pandas as pd
from tempfile import NamedTemporaryFile
import xlrd
import nltk
import numpy as np  # Keep this if used
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from docx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as PPTX_RT
import camelot
import easyocr
# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="Enhanced Document Parser for RAG", version="4.0.0")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

class CleaningOptions(BaseModel):
    normalize_unicode: bool = True
    remove_urls: bool = True
    remove_emails: bool = True
    clean_whitespace: bool = True
    preserve_structure: bool = True
    enable_ocr: bool = True
    aggressive_cleaning: bool = True

class ChunkingOptions(BaseModel):
    use_sentence_chunking: bool = True
    window_size: int = 8
    overlap: int = 2
    chunk_size: int = 1000
    chunk_overlap: int = 200

class DocumentChunk(BaseModel):
    content: str
    metadata: Dict[str, Any]
    chunk_id: str

class DocumentResult(BaseModel):
    text: str
    chunks: Optional[List[DocumentChunk]] = None
    metadata: Dict[str, Any]
    file_info: Dict[str, Any]
    processing_time: float
    fallback_used: bool = False
    ocr_used: bool = False

class EnhancedRAGTextCleaner:
    def __init__(self, options: CleaningOptions):
        self.options = options
    
    def clean_text(self, text: str, is_ocr: bool = False) -> str:
        """Enhanced text cleaning with OCR-specific fixes"""
        if not text:
            return ""
        
        if self.options.normalize_unicode:
            text = unicodedata.normalize('NFKC', text)
        
        if self.options.aggressive_cleaning:
            # Remove non-printable characters
            text = ''.join(c for c in text if c.isprintable() or c.isspace())
            
            # Fix common OCR ligatures and errors
            if is_ocr:
                text = text.replace('ﬁ', 'fi').replace('ﬂ', 'fl')
                text = re.sub(r'[|]', 'l', text)  # Common OCR error
                text = re.sub(r'(?<=[a-z])[0O](?=[a-z])', 'o', text)  # 0/O confusion
        
        if self.options.remove_urls:
            text = re.sub(r'http[s]?://\S+', '', text)
        
        if self.options.remove_emails:
            text = re.sub(r'\S+@\S+', '', text)
        
        if self.options.clean_whitespace:
            if self.options.preserve_structure:
                # Replace multiple spaces/tabs with single space
                text = re.sub(r'[ \t\r\f]+', ' ', text)
                # Normalize newlines but preserve paragraph breaks
                text = re.sub(r'\n+', '\n', text)
                # Remove empty lines
                text = '\n'.join([line.strip() for line in text.splitlines() if line.strip()])
            else:
                text = re.sub(r'\s+', ' ', text)
        
        # Remove excessive whitespace again after all processing
        text = re.sub(r' +', ' ', text)
        
        return text.strip()

class DocumentChunker:
    def __init__(self, options: ChunkingOptions):
        self.options = options
        # Download NLTK data if not present
        try:
            nltk.data.find('tokenizers/punkt')
        except LookupError:
            nltk.download('punkt')
    
    def create_chunks(self, content: str, file_path: str, metadata: Dict) -> List[DocumentChunk]:
        """Create document chunks with sentence-based or character-based chunking"""
        if self.options.use_sentence_chunking:
            return self._sentence_chunk(content, file_path, metadata)
        else:
            return self._character_chunk(content, file_path, metadata)
    
    def _sentence_chunk(self, content: str, file_path: str, metadata: Dict) -> List[DocumentChunk]:
        """Chunk text using sliding window over sentences for semantic coherence"""
        sentences = nltk.sent_tokenize(content)
        n = len(sentences)
        chunks = []
        chunk_id = 0
        i = 0
        
        while i < n:
            start = i
            end = min(i + self.options.window_size, n)
            chunk_text = ' '.join(sentences[start:end])
            
            if len(chunk_text.strip()) > 30:
                chunk_metadata = {
                    **metadata,
                    'chunk_id': chunk_id,
                    'chunk_type': 'sentence_based',
                    'sentence_start': start,
                    'sentence_end': end,
                    'total_sentences': n
                }
                
                chunks.append(DocumentChunk(
                    content=chunk_text.strip(),
                    metadata=chunk_metadata,
                    chunk_id=f"{os.path.splitext(os.path.basename(file_path))[0]}_{chunk_id}"
                ))
                chunk_id += 1
            
            i += self.options.window_size - self.options.overlap
        
        return chunks
    
    def _character_chunk(self, content: str, file_path: str, metadata: Dict) -> List[DocumentChunk]:
        """Traditional character-based chunking with overlap"""
        chunks = []
        chunk_id = 0
        start = 0
        
        while start < len(content):
            end = start + self.options.chunk_size
            chunk_text = content[start:end]
            
            # Try to break at word boundary if not at end
            if end < len(content):
                last_space = chunk_text.rfind(' ')
                if last_space > start + self.options.chunk_size // 2:
                    chunk_text = chunk_text[:last_space]
                    end = start + last_space
            
            if chunk_text.strip():
                chunk_metadata = {
                    **metadata,
                    'chunk_id': chunk_id,
                    'chunk_type': 'character_based',
                    'char_start': start,
                    'char_end': end
                }
                
                chunks.append(DocumentChunk(
                    content=chunk_text.strip(),
                    metadata=chunk_metadata,
                    chunk_id=f"{os.path.splitext(os.path.basename(file_path))[0]}_{chunk_id}"
                ))
                chunk_id += 1
            
            start = max(end - self.options.chunk_overlap, start + 1)
        
        return chunks

def merge_and_deduplicate_texts(texts: list[str]) -> str:
    """Merge multiple text sources, deduplicate lines, and preserve order."""
    seen = set()
    merged = []
    for text in texts:
        for line in text.splitlines():
            line = line.strip()
            if line and line not in seen:
                seen.add(line)
                merged.append(line)
    return '\n'.join(merged)


def extract_pdf(file_bytes: bytes, enable_ocr: bool = True, extract_tables: bool = True) -> tuple[str, dict]:
    """Optimized PDF extraction: combine PyMuPDF (layout-aware), PyPDF2, Camelot tables, and OCR for every page."""
    results = []
    meta = {"pages": 0, "methods": [], "ocr_pages": [], "fallback_used": False, "tables": 0, "table_pages": [], "layout_blocks": 0}
    
    # PyMuPDF layout-aware extraction
    try:
        print(f"Extracting with PyMuPDF layout-aware...")
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        meta["pages"] = doc.page_count
        print(f"Number of pages: {meta['pages']}")
        meta["methods"].append("PyMuPDF-layout")
        layout_texts = []
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            # Extract text blocks (layout-aware)
            blocks = page.get_text("blocks", sort=True)
            block_texts = [b[4] for b in blocks if b[6] == 0 and b[4].strip()]
            layout_texts.extend(block_texts)
        results.extend(layout_texts)
        meta["layout_blocks"] = len(layout_texts)
        doc.close()
    except Exception as e:
        logger.error(f"PyMuPDF layout-aware extraction failed: {e}")
        meta["fallback_used"] = True
    
    # PyMuPDF dict extraction (for structure)
    try:
        print(f"Extracting with PyMuPDF dict...")
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            d = page.get_text("dict", sort=True)
            for block in d.get("blocks", []):
                if block.get("type") == 0:  # text block
                    for line in block.get("lines", []):
                        for span in line.get("spans", []):
                            if span.get("text", "").strip():
                                results.append(span["text"])
        doc.close()
        meta["methods"].append("PyMuPDF-dict")
    except Exception as e:
        logger.error(f"PyMuPDF dict extraction failed: {e}")
        meta["fallback_used"] = True
    
    # PyPDF2
    try:
        print(f"Extracting with PyPDF2...")
        reader = PdfReader(BytesIO(file_bytes))
        meta["methods"].append("PyPDF2")
        pypdf2_texts = [page.extract_text() or "" for page in reader.pages]
        results.extend(pypdf2_texts)
    except Exception as e:
        logger.error(f"PyPDF2 failed: {e}")
        meta["fallback_used"] = True
    
    # Camelot table extraction
    if extract_tables:
        try:
            print("Extracting tables with Camelot...")
            with NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name

            table_texts = []
            tables = None

            try:
                tables = camelot.read_pdf(tmp_path, pages="all", flavor="lattice")
                if tables.n == 0:
                    raise ValueError("No tables found with lattice")
            except Exception as e:
                logger.warning(f"Lattice mode failed or returned no tables: {e}. Trying stream mode.")
                try:
                    tables = camelot.read_pdf(tmp_path, pages="all", flavor="stream")
                except Exception as e_stream:
                    logger.warning(f"Stream mode also failed: {e_stream}")
                    tables = None

            if tables:
                for table in tables:
                    df = table.df
                    table_str = df.to_string(index=False, header=True)
                    table_texts.append(f"[Table page {table.parsing_report.get('page', '?')}]\n{table_str}")
                    meta["table_pages"].append(table.parsing_report.get("page", "?"))

            if table_texts:
                results.extend(table_texts)
                meta["tables"] = len(table_texts)
                meta["methods"].append("Camelot")
            else:
                meta["tables"] = 0
                meta["table_pages"] = []
                meta["methods"].append("Camelot-no-tables")

        except Exception as e:
            logger.warning(f"Camelot table extraction failed: {e}")
            meta["tables"] = 0
            meta["table_pages"] = []
            meta["methods"].append("Camelot-error")
        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
    else:
        meta["tables"] = 0
        meta["table_pages"] = []
        meta["methods"].append("Camelot-skipped")

    
    # OCR for every page
    ocr_texts = []
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            pix = page.get_pixmap()
            img_data = pix.tobytes("png")
            img = Image.open(BytesIO(img_data))
            ocr_text = pytesseract.image_to_string(img)
            if ocr_text.strip():
                ocr_texts.append(ocr_text)
                meta["ocr_pages"].append(page_num)
        doc.close()
        if ocr_texts:
            results.extend(ocr_texts)
            meta["methods"].append("OCR-Tesseract")
    except Exception as e:
        logger.warning(f"OCR failed for PDF pages: {e}")
    
    # Merge and deduplicate
    merged = merge_and_deduplicate_texts(results)
    return merged.strip(), meta


def extract_image(file_bytes: bytes) -> tuple[str, dict]:
    """Optimized image extraction: run Tesseract and EasyOCR, merge results."""
    results = []
    meta = {"methods": [], "preprocessing": [], "fallback_used": False}
    
    # Tesseract OCR (basic and preprocessed)
    print(f"Extracting with Tesseract...")
    try:
        img = Image.open(BytesIO(file_bytes))
        if img.mode != 'RGB':
            img = img.convert('RGB')
        text_basic = pytesseract.image_to_string(img)
        results.append(text_basic)
        meta["methods"].append("Tesseract-basic")
    except Exception as e:
        logger.warning(f"Tesseract basic OCR failed: {e}")
        meta["fallback_used"] = True
    
    # Tesseract OCR (preprocessed)
    print(f"Extracting with Tesseract preprocessed...")
    try:
        img = Image.open(BytesIO(file_bytes))
        if img.mode != 'RGB':
            img = img.convert('RGB')
        img_array = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
        gray = cv2.cvtColor(img_array, cv2.COLOR_BGR2GRAY)
        gray = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1]
        custom_config = r'--oem 3 --psm 6'
        text_preproc = pytesseract.image_to_string(gray, config=custom_config)
        results.append(text_preproc)
        meta["methods"].append("Tesseract-preprocessed")
        meta["preprocessing"].append("OTSU_threshold")
    except Exception as e:
        logger.warning(f"Tesseract preprocessed OCR failed: {e}")
        meta["fallback_used"] = True
    
    # EasyOCR
    print(f"Extracting with EasyOCR...")
    try:
        reader = easyocr.Reader(['en'], gpu=False)
        img = np.array(Image.open(BytesIO(file_bytes)))
        easyocr_results = reader.readtext(img, detail=0)
        if easyocr_results:
            results.extend(easyocr_results)
            meta["methods"].append("EasyOCR")
    except Exception as e:
        logger.warning(f"EasyOCR failed: {e}")
        meta["fallback_used"] = True
    
    # Merge and deduplicate
    merged = merge_and_deduplicate_texts(results)
    meta["image_size"] = img.shape if 'img' in locals() else None
    return merged.strip(), meta

def extract_email_text(file_bytes: bytes) -> tuple[str, dict]:
    """Extract content from email files (.eml)"""
    try:
        content_str = file_bytes.decode('utf-8', errors='ignore')
        msg = email.message_from_string(content_str)
        
        text_content = f"Subject: {msg.get('Subject', '')}\n"
        text_content += f"From: {msg.get('From', '')}\n"
        text_content += f"To: {msg.get('To', '')}\n"
        text_content += f"Date: {msg.get('Date', '')}\n\n"
        
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        text_content += payload.decode('utf-8', errors='ignore')
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                text_content += payload.decode('utf-8', errors='ignore')
        
        return text_content.strip(), {
            "subject": msg.get('Subject', ''),
            "from": msg.get('From', ''),
            "multipart": msg.is_multipart()
        }
    except Exception as e:
        logger.error(f"Error reading email: {e}")
        return "", {"error": str(e)}

def extract_docx(file_bytes: bytes) -> tuple[str, dict]:
    """Optimized DOCX extraction: visible text, tables, comments, and core properties."""
    results = []
    meta = {"methods": [], "comments": 0, "tables": 0, "core_properties": {}, "mode": "optimized"}
    try:
        doc = Document(BytesIO(file_bytes))
        # Visible text
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        results.extend(paragraphs)
        meta["methods"].append("paragraphs")
        # Tables
        table_text = []
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                table_text.append(" | ".join(row_text))
        if table_text:
            results.extend(table_text)
            meta["tables"] = len(doc.tables)
            meta["methods"].append("tables")
        # Comments (if any)
        try:
            comments = []
            if hasattr(doc, 'part') and hasattr(doc.part, 'package'):
                for rel in doc.part.rels.values():
                    if rel.reltype == RT.COMMENTS:
                        comments_part = rel.target_part
                        for c in comments_part.element.findall('.//w:comment', namespaces=comments_part.element.nsmap):
                            text = ''.join([t.text for t in c.findall('.//w:t', namespaces=comments_part.element.nsmap) if t.text])
                            if text:
                                comments.append(text)
            if comments:
                results.extend(comments)
                meta["comments"] = len(comments)
                meta["methods"].append("comments")
        except Exception as e:
            logger.warning(f"DOCX comments extraction failed: {e}")
        # Core properties/metadata
        try:
            core_props = doc.core_properties
            core_dict = {k: getattr(core_props, k) for k in dir(core_props) if not k.startswith('_') and isinstance(getattr(core_props, k), (str, int, float, type(None)))}
            meta["core_properties"] = core_dict
            meta["methods"].append("core_properties")
        except Exception as e:
            logger.warning(f"DOCX core properties extraction failed: {e}")
        merged = merge_and_deduplicate_texts(results)
        return merged.strip(), meta
    except Exception as e:
        logger.error(f"Error reading DOCX: {e}")
        return "", {"error": str(e)}


def extract_pptx(file_bytes: bytes) -> tuple[str, dict]:
    """Optimized PPTX extraction: slides, notes, shapes, and core properties."""
    results = []
    meta = {"methods": [], "slides": 0, "notes": 0, "core_properties": {}, "mode": "optimized"}
    try:
        prs = Presentation(BytesIO(file_bytes))
        # Slides text
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            if slide_content:
                slides_text.append(f"Slide {i}:\n" + "\n".join(slide_content))
        if slides_text:
            results.extend(slides_text)
            meta["slides"] = len(prs.slides)
            meta["methods"].append("slides")
        # Notes
        notes_text = []
        for i, slide in enumerate(prs.slides, 1):
            if hasattr(slide, 'has_notes_slide') and slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes and notes.strip():
                    notes_text.append(f"Notes for Slide {i}:\n{notes.strip()}")
        if notes_text:
            results.extend(notes_text)
            meta["notes"] = len(notes_text)
            meta["methods"].append("notes")
        # Core properties/metadata
        try:
            core_props = prs.core_properties
            core_dict = {k: getattr(core_props, k) for k in dir(core_props) if not k.startswith('_') and isinstance(getattr(core_props, k), (str, int, float, type(None)))}
            meta["core_properties"] = core_dict
            meta["methods"].append("core_properties")
        except Exception as e:
            logger.warning(f"PPTX core properties extraction failed: {e}")
        merged = merge_and_deduplicate_texts(results)
        return merged.strip(), meta
    except Exception as e:
        logger.error(f"Error reading PPTX: {e}")
        return "", {"error": str(e)}

def extract_excel_text(file_bytes: bytes, ext: str) -> tuple[str, dict]:
    with NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        if ext == ".xlsx":
            df_dict = pd.read_excel(tmp_path, sheet_name=None)
            sheets_text = [f"Sheet: {name}\n{df.to_string(index=False)}" 
                          for name, df in df_dict.items()]
            return "\n\n".join(sheets_text), {"sheets": len(df_dict)}
        else:
            book = xlrd.open_workbook(tmp_path)
            sheets_text = []
            for sheet in book.sheets():
                sheet_data = []
                for row_idx in range(sheet.nrows):
                    row = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                    sheet_data.append(", ".join(row))
                sheets_text.append(f"Sheet: {sheet.name}\n" + "\n".join(sheet_data))
            return "\n\n".join(sheets_text), {"sheets": len(book.sheets())}
    finally:
        os.remove(tmp_path)

EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".png": extract_image,
    ".jpg": extract_image,
    ".jpeg": extract_image,
    ".tiff": extract_image,
    ".bmp": extract_image,
    ".txt": lambda content, **kwargs: (content.decode("utf-8", errors="ignore").strip(), {"encoding": "utf-8"}),
    ".eml": lambda content, **kwargs: extract_email_text(content),
    ".html": lambda content, **kwargs: (BeautifulSoup(content, "html.parser").get_text(separator='\n').strip(), {"parser": "html.parser"}),
    ".csv": lambda content, **kwargs: pd.read_csv(BytesIO(content)).to_string(index=False) if content else ("", {}),
    ".json": lambda content, **kwargs: (json.dumps(json.loads(content.decode("utf-8", errors="ignore")), indent=2, ensure_ascii=False), {"valid": True}),
    ".xlsx": lambda content, ext, **kwargs: extract_excel_text(content, ext),
    ".xls": lambda content, ext, **kwargs: extract_excel_text(content, ext),
}

@app.get("/")
async def root():
    return {
        "message": "Enhanced Document Parser for RAG", 
        "version": "4.0.0",
        "status": "running",
        "supported_formats": list(EXTRACTORS.keys()),
        "features": ["OCR support", "Sentence-based chunking", "Fallback mechanisms", "Enhanced cleaning"],
        "endpoints": {
            "parse": "/parse - Parse single document",
            "parse_batch": "/parse-batch - Parse multiple documents",
            "health": "/health - Health check"
        }
    }
@app.get("/health")
async def health_check():
    return {"status": "healthy", "timestamp": datetime.now().isoformat()}
@app.post("/parse", response_model=DocumentResult)
async def parse_file(
    file: UploadFile = File(...),
    # Cleaning options
    normalize_unicode: bool = Query(True),
    remove_urls: bool = Query(True),
    remove_emails: bool = Query(True),
    clean_whitespace: bool = Query(True),
    preserve_structure: bool = Query(True),
    enable_ocr: bool = Query(True),
    aggressive_cleaning: bool = Query(True),
    # Chunking options
    create_chunks: bool = Query(False),
    use_sentence_chunking: bool = Query(True),
    window_size: int = Query(8),
    overlap: int = Query(2),
    chunk_size: int = Query(1000),
    chunk_overlap: int = Query(200),
    # Extraction options
    extract_tables: bool = Query(True, description="Extract tables using Camelot (may be slow for large PDFs).")
):
    start_time = datetime.now()
    fallback_used = False
    ocr_used = False
    
    try:
        content = await file.read()
        ext = os.path.splitext(file.filename)[-1].lower()
        
        if ext not in EXTRACTORS:
            raise HTTPException(400, f"Unsupported file type: {ext}. Supported: {list(EXTRACTORS.keys())}")
        
        # Extract text with enhanced methods
        extractor_kwargs = {'enable_ocr': enable_ocr, 'extract_tables': extract_tables}
        
        # Select only relevant kwargs for the extractor
        sig = inspect.signature(EXTRACTORS[ext])
        relevant_kwargs = {k: v for k, v in extractor_kwargs.items() if k in sig.parameters}
        
        if 'content' in sig.parameters:
            raw_text, metadata = EXTRACTORS[ext](content, **relevant_kwargs)
        else: # Handle callables that might not be wrapped to accept **kwargs
            if ext in [".xlsx", ".xls"]:
                raw_text, metadata = EXTRACTORS[ext](content, ext)
            elif ext == ".pdf":
                 raw_text, metadata = EXTRACTORS[ext](content, enable_ocr=enable_ocr, extract_tables=extract_tables)
            else:
                 raw_text, metadata = EXTRACTORS[ext](content)

        ocr_used = metadata.get("ocr_used", False) or (isinstance(metadata.get("ocr_pages"), list) and len(metadata.get("ocr_pages")) > 0)
        fallback_used = metadata.get("fallback_used", False)
        
        # Enhanced cleaning
        cleaning_options = CleaningOptions(
            normalize_unicode=normalize_unicode,
            remove_urls=remove_urls,
            remove_emails=remove_emails,
            clean_whitespace=clean_whitespace,
            preserve_structure=preserve_structure,
            enable_ocr=enable_ocr,
            aggressive_cleaning=aggressive_cleaning
        )
        
        cleaner = EnhancedRAGTextCleaner(cleaning_options)
        cleaned_text = cleaner.clean_text(raw_text, is_ocr=ocr_used)
        
        # Create chunks if requested
        chunks = None
        if create_chunks and cleaned_text.strip():
            chunking_options = ChunkingOptions(
                use_sentence_chunking=use_sentence_chunking,
                window_size=window_size,
                overlap=overlap,
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap
            )
            
            chunker = DocumentChunker(chunking_options)
            enhanced_metadata = {
                **metadata,
                'source': file.filename,
                'file_type': ext,
                'ocr_used': ocr_used,
                'fallback_used': fallback_used,
                'cleaning_steps': 'unicode_normalize,remove_nonprintable,whitespace,ocr_fixes'
            }
            chunks = chunker.create_chunks(cleaned_text, file.filename, enhanced_metadata)
        
        processing_time = (datetime.now() - start_time).total_seconds()
        
        return DocumentResult(
            text=cleaned_text,
            chunks=chunks,
            metadata={
                **metadata,
                'processing_steps': 'extract,clean,chunk' if chunks else 'extract,clean',
                'total_chunks': len(chunks) if chunks else 0
            },
            processing_time=processing_time,
            file_info={
                "filename": file.filename,
                "file_size": len(content),
                "file_type": ext
            },
            fallback_used=fallback_used,
            ocr_used=ocr_used
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Processing failed: {str(e)}")
        raise HTTPException(500, f"Processing failed: {str(e)}")

@app.post("/parse-batch")
async def parse_batch(
    files: List[UploadFile] = File(...),
    max_workers: int = Query(4, ge=1, le=10),
    # Cleaning options
    normalize_unicode: bool = Query(True),
    remove_urls: bool = Query(True),
    remove_emails: bool = Query(True),
    clean_whitespace: bool = Query(True),
    preserve_structure: bool = Query(True),
    enable_ocr: bool = Query(True),
    aggressive_cleaning: bool = Query(True),
    # Chunking options
    create_chunks: bool = Query(False),
    use_sentence_chunking: bool = Query(True),
    window_size: int = Query(8),
    overlap: int = Query(2),
    chunk_size: int = Query(1000),
    chunk_overlap: int = Query(200),
    # Extraction options
    extract_tables: bool = Query(True, description="Extract tables using Camelot (may be slow for large PDFs).")
):
    """Process multiple files in parallel"""
    start_time = datetime.now()
    results = []
    
    async def process_single_file(file: UploadFile):
        try:
            # Read file content
            content = await file.read()
            ext = os.path.splitext(file.filename)[-1].lower()
            
            if ext not in EXTRACTORS:
                return {"filename": file.filename, "error": f"Unsupported file type: {ext}"}
            
            # Process with default settings
            if ext == ".pdf":
                raw_text, metadata = extract_pdf(content, enable_ocr=True, extract_tables=extract_tables)
            elif ext in [".xlsx", ".xls"]:
                raw_text, metadata = extract_excel_text(content, ext)
            elif ext in [".docx"]:
                 raw_text, metadata = extract_docx(content)
            elif ext in [".pptx"]:
                 raw_text, metadata = extract_pptx(content)
            elif ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
                 raw_text, metadata = extract_image(content)
            else:
                raw_text, metadata = EXTRACTORS[ext](content)
            
            # Basic cleaning
            cleaner = EnhancedRAGTextCleaner(CleaningOptions(
                normalize_unicode=normalize_unicode,
                remove_urls=remove_urls,
                remove_emails=remove_emails,
                clean_whitespace=clean_whitespace,
                preserve_structure=preserve_structure,
                enable_ocr=enable_ocr,
                aggressive_cleaning=aggressive_cleaning
            ))
            cleaned_text = cleaner.clean_text(raw_text)
            
            return {
                "filename": file.filename,
                "text": cleaned_text,
                "metadata": metadata,
                "success": True
            }
            
        except Exception as e:
            logger.error(f"Error processing {file.filename}: {e}")
            return {"filename": file.filename, "error": str(e), "success": False}
    
    # Process files concurrently
    tasks = [process_single_file(file) for file in files]
    results = await asyncio.gather(*tasks, return_exceptions=True)
    
    processing_time = (datetime.now() - start_time).total_seconds()
    
    return {
        "results": results,
        "total_files": len(files),
        "processing_time": processing_time,
        "successful": len([r for r in results if isinstance(r, dict) and r.get("success")])
    }

if __name__ == "__main__":
    import uvicorn
    import os
    import inspect

    port = int(os.getenv("PORT", 8001))  # Default to 8000 if not set
    uvicorn.run(app, host="0.0.0.0", port=port)
