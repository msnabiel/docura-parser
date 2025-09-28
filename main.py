from fastapi import FastAPI, UploadFile, File, HTTPException, Query
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Dict, Any, List, Optional
from io import BytesIO
import os, re, unicodedata, json, csv, email, logging
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed
import asyncio
from functools import partial

# Core extraction libraries
from PyPDF2 import PdfReader
import fitz
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import cv2
from bs4 import BeautifulSoup
import pandas as pd
from tempfile import NamedTemporaryFile
import xlrd
import nltk
import numpy as np  # Keep this if used
# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="Enhanced Document Parser for RAG", version="4.0.0")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

class CleaningOptions(BaseModel):
    normalize_unicode: bool = True
    remove_urls: bool = False
    remove_emails: bool = False
    clean_whitespace: bool = True
    preserve_structure: bool = True
    max_length: int = 100000
    enable_ocr: bool = True
    aggressive_cleaning: bool = True

class ChunkingOptions(BaseModel):
    use_sentence_chunking: bool = True
    window_size: int = 8
    overlap: int = 2
    chunk_size: int = 1000
    chunk_overlap: int = 200

class DocumentChunk(BaseModel):
    content: str
    metadata: Dict[str, Any]
    chunk_id: str

class DocumentResult(BaseModel):
    text: str
    chunks: Optional[List[DocumentChunk]] = None
    metadata: Dict[str, Any]
    file_info: Dict[str, Any]
    processing_time: float
    fallback_used: bool = False
    ocr_used: bool = False

class EnhancedRAGTextCleaner:
    def __init__(self, options: CleaningOptions):
        self.options = options
    
    def clean_text(self, text: str, is_ocr: bool = False) -> str:
        """Enhanced text cleaning with OCR-specific fixes"""
        if not text:
            return ""
        
        if self.options.normalize_unicode:
            text = unicodedata.normalize('NFKC', text)
        
        if self.options.aggressive_cleaning:
            # Remove non-printable characters
            text = ''.join(c for c in text if c.isprintable() or c.isspace())
            
            # Fix common OCR ligatures and errors
            if is_ocr:
                text = text.replace('ﬁ', 'fi').replace('ﬂ', 'fl')
                text = re.sub(r'[|]', 'l', text)  # Common OCR error
                text = re.sub(r'(?<=[a-z])[0O](?=[a-z])', 'o', text)  # 0/O confusion
        
        if self.options.remove_urls:
            text = re.sub(r'http[s]?://\S+', '', text)
        
        if self.options.remove_emails:
            text = re.sub(r'\S+@\S+', '', text)
        
        if self.options.clean_whitespace:
            if self.options.preserve_structure:
                # Replace multiple spaces/tabs with single space
                text = re.sub(r'[ \t\r\f]+', ' ', text)
                # Normalize newlines but preserve paragraph breaks
                text = re.sub(r'\n+', '\n', text)
                # Remove empty lines
                text = '\n'.join([line.strip() for line in text.splitlines() if line.strip()])
            else:
                text = re.sub(r'\s+', ' ', text)
        
        # Remove excessive whitespace again after all processing
        text = re.sub(r' +', ' ', text)
        
        if len(text) > self.options.max_length:
            text = text[:self.options.max_length].rsplit(' ', 1)[0] + "..."
        
        return text.strip()

class DocumentChunker:
    def __init__(self, options: ChunkingOptions):
        self.options = options
        # Download NLTK data if not present
        try:
            nltk.data.find('tokenizers/punkt')
        except LookupError:
            nltk.download('punkt')
    
    def create_chunks(self, content: str, file_path: str, metadata: Dict) -> List[DocumentChunk]:
        """Create document chunks with sentence-based or character-based chunking"""
        if self.options.use_sentence_chunking:
            return self._sentence_chunk(content, file_path, metadata)
        else:
            return self._character_chunk(content, file_path, metadata)
    
    def _sentence_chunk(self, content: str, file_path: str, metadata: Dict) -> List[DocumentChunk]:
        """Chunk text using sliding window over sentences for semantic coherence"""
        sentences = nltk.sent_tokenize(content)
        n = len(sentences)
        chunks = []
        chunk_id = 0
        i = 0
        
        while i < n:
            start = i
            end = min(i + self.options.window_size, n)
            chunk_text = ' '.join(sentences[start:end])
            
            if len(chunk_text.strip()) > 30:
                chunk_metadata = {
                    **metadata,
                    'chunk_id': chunk_id,
                    'chunk_type': 'sentence_based',
                    'sentence_start': start,
                    'sentence_end': end,
                    'total_sentences': n
                }
                
                chunks.append(DocumentChunk(
                    content=chunk_text.strip(),
                    metadata=chunk_metadata,
                    chunk_id=f"{os.path.splitext(os.path.basename(file_path))[0]}_{chunk_id}"
                ))
                chunk_id += 1
            
            i += self.options.window_size - self.options.overlap
        
        return chunks
    
    def _character_chunk(self, content: str, file_path: str, metadata: Dict) -> List[DocumentChunk]:
        """Traditional character-based chunking with overlap"""
        chunks = []
        chunk_id = 0
        start = 0
        
        while start < len(content):
            end = start + self.options.chunk_size
            chunk_text = content[start:end]
            
            # Try to break at word boundary if not at end
            if end < len(content):
                last_space = chunk_text.rfind(' ')
                if last_space > start + self.options.chunk_size // 2:
                    chunk_text = chunk_text[:last_space]
                    end = start + last_space
            
            if chunk_text.strip():
                chunk_metadata = {
                    **metadata,
                    'chunk_id': chunk_id,
                    'chunk_type': 'character_based',
                    'char_start': start,
                    'char_end': end
                }
                
                chunks.append(DocumentChunk(
                    content=chunk_text.strip(),
                    metadata=chunk_metadata,
                    chunk_id=f"{os.path.splitext(os.path.basename(file_path))[0]}_{chunk_id}"
                ))
                chunk_id += 1
            
            start = max(end - self.options.chunk_overlap, start + 1)
        
        return chunks

def extract_pdf_text_with_ocr(file_bytes: bytes, enable_ocr: bool = True) -> tuple[str, dict]:
    """Enhanced PDF extraction with OCR fallback"""
    text = ""
    metadata = {"pages": 0, "method": "unknown", "ocr_used": False}
    fallback_used = False
    
    try:
        # First try PyMuPDF for better text extraction
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        metadata["pages"] = doc.page_count
        metadata["method"] = "PyMuPDF"
        
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            page_text = page.get_text()
            
            # If no text found and OCR is enabled, try OCR
            if not page_text.strip() and enable_ocr:
                try:
                    # Convert page to image
                    pix = page.get_pixmap()
                    img_data = pix.tobytes("png")
                    img = Image.open(BytesIO(img_data))
                    
                    # OCR the image
                    ocr_text = pytesseract.image_to_string(img)
                    page_text = ocr_text
                    metadata["ocr_used"] = True
                    logger.info(f"Used OCR for page {page_num + 1}")
                except Exception as e:
                    logger.warning(f"OCR failed for page {page_num + 1}: {e}")
            
            text += page_text + "\n"
        
        doc.close()
        
    except Exception as e:
        logger.error(f"Error with PyMuPDF, falling back to PyPDF2: {e}")
        fallback_used = True
        try:
            reader = PdfReader(BytesIO(file_bytes))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            metadata = {"pages": len(reader.pages), "method": "PyPDF2_fallback"}
        except Exception as e2:
            logger.error(f"Error with PyPDF2 fallback: {e2}")
            raise HTTPException(500, f"Failed to extract PDF text: {str(e2)}")
    
    metadata["fallback_used"] = fallback_used
    return text.strip(), metadata

def extract_image_text_enhanced(file_bytes: bytes) -> tuple[str, dict]:
    """Enhanced image text extraction with preprocessing"""
    try:
        # Load image
        img = Image.open(BytesIO(file_bytes))
        if img.mode != 'RGB':
            img = img.convert('RGB')
        
        # Convert PIL to OpenCV format for preprocessing
        img_array = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
        gray = cv2.cvtColor(img_array, cv2.COLOR_BGR2GRAY)
        
        # Apply preprocessing to improve OCR
        gray = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1]
        
        # OCR with enhanced config
        custom_config = r'--oem 3 --psm 6'
        text = pytesseract.image_to_string(gray, config=custom_config)
        
        return text.strip(), {"image_size": img.size, "preprocessing": "OTSU_threshold"}
    
    except Exception as e:
        # Fallback to basic OCR without preprocessing
        logger.warning(f"Enhanced OCR failed, falling back to basic: {e}")
        img = Image.open(BytesIO(file_bytes))
        if img.mode != 'RGB':
            img = img.convert('RGB')
        text = pytesseract.image_to_string(img)
        return text.strip(), {"image_size": img.size, "fallback_used": True}

def extract_email_text(file_bytes: bytes) -> tuple[str, dict]:
    """Extract content from email files (.eml)"""
    try:
        content_str = file_bytes.decode('utf-8', errors='ignore')
        msg = email.message_from_string(content_str)
        
        text_content = f"Subject: {msg.get('Subject', '')}\n"
        text_content += f"From: {msg.get('From', '')}\n"
        text_content += f"To: {msg.get('To', '')}\n"
        text_content += f"Date: {msg.get('Date', '')}\n\n"
        
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        text_content += payload.decode('utf-8', errors='ignore')
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                text_content += payload.decode('utf-8', errors='ignore')
        
        return text_content.strip(), {
            "subject": msg.get('Subject', ''),
            "from": msg.get('From', ''),
            "multipart": msg.is_multipart()
        }
    except Exception as e:
        logger.error(f"Error reading email: {e}")
        return "", {"error": str(e)}

def extract_docx_text(file_bytes: bytes) -> tuple[str, dict]:
    """Enhanced DOCX extraction"""
    try:
        doc = Document(BytesIO(file_bytes))
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        
        # Extract table content
        table_text = []
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                table_text.append(" | ".join(row_text))
        
        text = "\n".join(paragraphs)
        if table_text:
            text += "\n\nTables:\n" + "\n".join(table_text)
        
        return text.strip(), {
            "paragraphs": len(paragraphs), 
            "tables": len(doc.tables),
            "method": "python-docx"
        }
    except Exception as e:
        logger.error(f"Error reading DOCX: {e}")
        return "", {"error": str(e)}

def extract_pptx_text(file_bytes: bytes) -> tuple[str, dict]:
    prs = Presentation(BytesIO(file_bytes))
    slides_text = []
    
    for i, slide in enumerate(prs.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())
        if slide_content:
            slides_text.append(f"Slide {i}:\n" + "\n".join(slide_content))
    
    return "\n\n".join(slides_text), {"slides": len(prs.slides)}

# Update extractors to include new formats and enhanced versions
EXTRACTORS = {
    ".pdf": extract_pdf_text_with_ocr,
    ".docx": extract_docx_text,
    ".pptx": extract_pptx_text,
    ".png": extract_image_text_enhanced,
    ".jpg": extract_image_text_enhanced,
    ".jpeg": extract_image_text_enhanced,
    ".tiff": extract_image_text_enhanced,
    ".bmp": extract_image_text_enhanced,
    ".txt": lambda content: (content.decode("utf-8", errors="ignore").strip(), {"encoding": "utf-8"}),
    ".eml": extract_email_text,
    ".html": lambda content: (BeautifulSoup(content, "html.parser").get_text(separator='\n').strip(), {"parser": "html.parser"}),
    ".csv": lambda content: pd.read_csv(BytesIO(content)).to_string(index=False) if content else ("", {}),
    ".json": lambda content: (json.dumps(json.loads(content.decode("utf-8", errors="ignore")), indent=2, ensure_ascii=False), {"valid": True}),
    ".xlsx": lambda content, ext: extract_excel_text(content, ext),
    ".xls": lambda content, ext: extract_excel_text(content, ext),
}

def extract_excel_text(file_bytes: bytes, ext: str) -> tuple[str, dict]:
    with NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        if ext == ".xlsx":
            df_dict = pd.read_excel(tmp_path, sheet_name=None)
            sheets_text = [f"Sheet: {name}\n{df.to_string(index=False)}" 
                          for name, df in df_dict.items()]
            return "\n\n".join(sheets_text), {"sheets": len(df_dict)}
        else:
            book = xlrd.open_workbook(tmp_path)
            sheets_text = []
            for sheet in book.sheets():
                sheet_data = []
                for row_idx in range(sheet.nrows):
                    row = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                    sheet_data.append(", ".join(row))
                sheets_text.append(f"Sheet: {sheet.name}\n" + "\n".join(sheet_data))
            return "\n\n".join(sheets_text), {"sheets": len(book.sheets())}
    finally:
        os.remove(tmp_path)

@app.get("/")
async def root():
    return {
        "message": "Enhanced Document Parser for RAG", 
        "version": "4.0.0",
        "status": "running",
        "supported_formats": list(EXTRACTORS.keys()),
        "features": ["OCR support", "Sentence-based chunking", "Fallback mechanisms", "Enhanced cleaning"],
        "endpoints": {
            "parse": "/parse - Parse single document",
            "parse_batch": "/parse-batch - Parse multiple documents",
            "health": "/health - Health check"
        }
    }
@app.get("/health")
async def health_check():
    return {"status": "healthy", "timestamp": datetime.now().isoformat()}
@app.post("/parse", response_model=DocumentResult)
async def parse_file(
    file: UploadFile = File(...),
    # Cleaning options
    normalize_unicode: bool = Query(True),
    remove_urls: bool = Query(False),
    remove_emails: bool = Query(False),
    clean_whitespace: bool = Query(True),
    preserve_structure: bool = Query(True),
    max_length: int = Query(100000),
    enable_ocr: bool = Query(True),
    aggressive_cleaning: bool = Query(True),
    # Chunking options
    create_chunks: bool = Query(False),
    use_sentence_chunking: bool = Query(True),
    window_size: int = Query(8),
    overlap: int = Query(2),
    chunk_size: int = Query(1000),
    chunk_overlap: int = Query(200)
):
    start_time = datetime.now()
    fallback_used = False
    ocr_used = False
    
    try:
        content = await file.read()
        ext = os.path.splitext(file.filename)[-1].lower()
        
        if ext not in EXTRACTORS:
            raise HTTPException(400, f"Unsupported file type: {ext}. Supported: {list(EXTRACTORS.keys())}")
        
        # Extract text with enhanced methods
        if ext == ".pdf":
            raw_text, metadata = EXTRACTORS[ext](content, enable_ocr)
            ocr_used = metadata.get("ocr_used", False)
            fallback_used = metadata.get("fallback_used", False)
        elif ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
            raw_text, metadata = EXTRACTORS[ext](content)
            ocr_used = True
            fallback_used = metadata.get("fallback_used", False)
        elif ext in [".xlsx", ".xls"]:
            raw_text, metadata = EXTRACTORS[ext](content, ext)
        else:
            raw_text, metadata = EXTRACTORS[ext](content)
        
        # Enhanced cleaning
        cleaning_options = CleaningOptions(
            normalize_unicode=normalize_unicode,
            remove_urls=remove_urls,
            remove_emails=remove_emails,
            clean_whitespace=clean_whitespace,
            preserve_structure=preserve_structure,
            max_length=max_length,
            enable_ocr=enable_ocr,
            aggressive_cleaning=aggressive_cleaning
        )
        
        cleaner = EnhancedRAGTextCleaner(cleaning_options)
        cleaned_text = cleaner.clean_text(raw_text, is_ocr=ocr_used)
        
        # Create chunks if requested
        chunks = None
        if create_chunks and cleaned_text.strip():
            chunking_options = ChunkingOptions(
                use_sentence_chunking=use_sentence_chunking,
                window_size=window_size,
                overlap=overlap,
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap
            )
            
            chunker = DocumentChunker(chunking_options)
            enhanced_metadata = {
                **metadata,
                'source': file.filename,
                'file_type': ext,
                'ocr_used': ocr_used,
                'fallback_used': fallback_used,
                'cleaning_steps': 'unicode_normalize,remove_nonprintable,whitespace,ocr_fixes'
            }
            chunks = chunker.create_chunks(cleaned_text, file.filename, enhanced_metadata)
        
        processing_time = (datetime.now() - start_time).total_seconds()
        
        return DocumentResult(
            text=cleaned_text,
            chunks=chunks,
            metadata={
                **metadata,
                'processing_steps': 'extract,clean,chunk' if chunks else 'extract,clean',
                'total_chunks': len(chunks) if chunks else 0
            },
            processing_time=processing_time,
            file_info={
                "filename": file.filename,
                "file_size": len(content),
                "file_type": ext
            },
            fallback_used=fallback_used,
            ocr_used=ocr_used
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Processing failed: {str(e)}")
        raise HTTPException(500, f"Processing failed: {str(e)}")

@app.post("/parse-batch")
async def parse_batch(
    files: List[UploadFile] = File(...),
    max_workers: int = Query(4, ge=1, le=10)
):
    """Process multiple files in parallel"""
    start_time = datetime.now()
    results = []
    
    async def process_single_file(file: UploadFile):
        try:
            # Read file content
            content = await file.read()
            ext = os.path.splitext(file.filename)[-1].lower()
            
            if ext not in EXTRACTORS:
                return {"filename": file.filename, "error": f"Unsupported file type: {ext}"}
            
            # Process with default settings
            if ext == ".pdf":
                raw_text, metadata = EXTRACTORS[ext](content, True)
            elif ext in [".xlsx", ".xls"]:
                raw_text, metadata = EXTRACTORS[ext](content, ext)
            else:
                raw_text, metadata = EXTRACTORS[ext](content)
            
            # Basic cleaning
            cleaner = EnhancedRAGTextCleaner(CleaningOptions())
            cleaned_text = cleaner.clean_text(raw_text)
            
            return {
                "filename": file.filename,
                "text": cleaned_text,
                "metadata": metadata,
                "success": True
            }
            
        except Exception as e:
            logger.error(f"Error processing {file.filename}: {e}")
            return {"filename": file.filename, "error": str(e), "success": False}
    
    # Process files concurrently
    tasks = [process_single_file(file) for file in files]
    results = await asyncio.gather(*tasks, return_exceptions=True)
    
    processing_time = (datetime.now() - start_time).total_seconds()
    
    return {
        "results": results,
        "total_files": len(files),
        "processing_time": processing_time,
        "successful": len([r for r in results if isinstance(r, dict) and r.get("success")])
    }

if __name__ == "__main__":
    import uvicorn
    import os

    port = int(os.getenv("PORT", 8000))  # Default to 8000 if not set
    uvicorn.run(app, host="0.0.0.0", port=port)
