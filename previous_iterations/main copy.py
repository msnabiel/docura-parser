from fastapi import FastAPI, UploadFile, File
from fastapi.responses import JSONResponse
from typing import Optional
from io import BytesIO
import os
import logging
import csv
import json

from PyPDF2 import PdfReader
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
from bs4 import BeautifulSoup
import markdown
import openpyxl
import xlrd
from ebooklib import epub
from striprtf.striprtf import rtf_to_text
from tempfile import NamedTemporaryFile

app = FastAPI()
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("parser")

# --- Text extractors ---

def extract_pdf_text(file_bytes: bytes) -> str:
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        text = "\n".join(page.get_text() for page in doc)
    except Exception as e:
        logger.warning("fitz failed, fallback to PyPDF2: %s", e)
        reader = PdfReader(BytesIO(file_bytes))
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
    return text.strip()

def extract_docx_text(file_bytes: bytes) -> str:
    doc = Document(BytesIO(file_bytes))
    return "\n".join([para.text for para in doc.paragraphs]).strip()

def extract_pptx_text(file_bytes: bytes) -> str:
    prs = Presentation(BytesIO(file_bytes))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text).strip()

def extract_image_text(file_bytes: bytes) -> str:
    image = Image.open(BytesIO(file_bytes)).convert("RGB")
    return pytesseract.image_to_string(image).strip()

def extract_txt_text(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="ignore").strip()

def extract_html_text(file_bytes: bytes) -> str:
    soup = BeautifulSoup(file_bytes, "html.parser")
    return soup.get_text().strip()

def extract_md_text(file_bytes: bytes) -> str:
    html = markdown.markdown(file_bytes.decode("utf-8", errors="ignore"))
    soup = BeautifulSoup(html, "html.parser")
    return soup.get_text().strip()

def extract_csv_text(file_bytes: bytes) -> str:
    text = file_bytes.decode("utf-8", errors="ignore")
    reader = csv.reader(text.splitlines())
    return "\n".join([", ".join(row) for row in reader])

def extract_excel_text(file_bytes: bytes, ext: str) -> str:
    with NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        if ext == ".xlsx":
            wb = openpyxl.load_workbook(tmp_path, data_only=True)
            sheets = wb.sheetnames
            return "\n".join(
                [", ".join([str(cell.value or "") for cell in row]) for sheet in sheets for row in wb[sheet].iter_rows()]
            )
        elif ext == ".xls":
            book = xlrd.open_workbook(tmp_path)
            result = []
            for sheet in book.sheets():
                for row_idx in range(sheet.nrows):
                    result.append(", ".join(str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)))
            return "\n".join(result)
    finally:
        os.remove(tmp_path)

def extract_json_text(file_bytes: bytes) -> str:
    try:
        data = json.loads(file_bytes.decode("utf-8", errors="ignore"))
        return json.dumps(data, indent=2)
    except json.JSONDecodeError:
        return "Invalid JSON file."

def extract_epub_text(file_bytes: bytes) -> str:
    book = epub.read_epub(BytesIO(file_bytes))
    text = []
    for item in book.get_items():
        if item.get_type() == epub.ITEM_DOCUMENT:
            soup = BeautifulSoup(item.get_content(), "html.parser")
            text.append(soup.get_text())
    return "\n".join(text).strip()

def extract_rtf_text(file_bytes: bytes) -> str:
    return rtf_to_text(file_bytes.decode("utf-8", errors="ignore")).strip()

# --- API Route ---

@app.post("/parse")
async def parse_file(file: UploadFile = File(...)):
    content = await file.read()
    ext = os.path.splitext(file.filename)[-1].lower()

    try:
        if ext == ".pdf":
            text = extract_pdf_text(content)
        elif ext == ".docx":
            text = extract_docx_text(content)
        elif ext == ".pptx":
            text = extract_pptx_text(content)
        elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"]:
            text = extract_image_text(content)
        elif ext == ".txt":
            text = extract_txt_text(content)
        elif ext in [".html", ".htm"]:
            text = extract_html_text(content)
        elif ext == ".md":
            text = extract_md_text(content)
        elif ext == ".csv":
            text = extract_csv_text(content)
        elif ext == ".json":
            text = extract_json_text(content)
        elif ext in [".xlsx", ".xls"]:
            text = extract_excel_text(content, ext)
        elif ext == ".epub":
            text = extract_epub_text(content)
        elif ext == ".rtf":
            text = extract_rtf_text(content)
        else:
            return JSONResponse({"error": f"Unsupported file type: {ext}"}, status_code=400)
        return {"text": text[:10000]}  # Trimmed to 10k chars
    except Exception as e:
        logger.exception("Failed to parse file")
        return JSONResponse({"error": str(e)}, status_code=500)
